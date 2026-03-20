"""
Analytix AI — Professional Export Service.
Generates polished PDF reports and PPTX presentations.
"""
import io
import pandas as pd
import numpy as np
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from datetime import datetime
from typing import Any

# PDF
from reportlab.lib.pagesizes import A4, letter
from reportlab.lib.units import inch, cm, mm
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Paragraph, Spacer, Image, Table, TableStyle,
    PageBreak, HRFlowable, KeepTogether
)
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_RIGHT, TA_JUSTIFY
from reportlab.graphics.shapes import Drawing, Rect, String
from reportlab.graphics import renderPDF

# PPTX
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

import logging
logger = logging.getLogger(__name__)


# ── Brand Colors ──────────────────────────────────────
BRAND_PRIMARY = "#6366f1"
BRAND_SECONDARY = "#8b5cf6"
BRAND_ACCENT = "#06b6d4"
BRAND_DARK = "#1e1b4b"
BRAND_TEXT = "#334155"
BRAND_MUTED = "#64748b"
BRAND_LIGHT_BG = "#f8fafc"
BRAND_BORDER = "#e2e8f0"

CHART_PALETTE = [
    "#6366f1", "#8b5cf6", "#06b6d4", "#10b981", "#f59e0b",
    "#ef4444", "#ec4899", "#14b8a6", "#f97316", "#3b82f6",
]


class ExportService:
    """Generates professional PDF reports and PPTX presentations."""

    # ══════════════════════════════════════════════════
    #  CHART RENDERING (shared by PDF & PPTX)
    # ══════════════════════════════════════════════════
    def generate_chart_image(self, chart_config: dict, width=8, height=5, style="light") -> bytes:
        """Render a chart as a high-quality PNG image."""
        fig, ax = plt.subplots(figsize=(width, height))

        if style == "dark":
            fig.patch.set_facecolor("#0f172a")
            ax.set_facecolor("#1e293b")
            text_color, grid_color, spine_color = "white", "#334155", "#475569"
            tick_color = "#94a3b8"
        else:
            fig.patch.set_facecolor("white")
            ax.set_facecolor("#fafafa")
            text_color, grid_color, spine_color = "#1e293b", "#e2e8f0", "#cbd5e1"
            tick_color = "#64748b"

        chart_type = chart_config.get("chart_type", "bar")
        if hasattr(chart_type, "value"):
            chart_type = chart_type.value
        data = chart_config.get("data", [])
        title = chart_config.get("title", "Chart")
        x_col = chart_config.get("x_column", "")
        y_col = chart_config.get("y_column", "")

        if not data:
            ax.text(0.5, 0.5, "No data", ha="center", va="center",
                    color=text_color, fontsize=14, transform=ax.transAxes)
        else:
            if chart_type == "bar":
                labels = [str(d.get(x_col, ""))[:20] for d in data]
                values = [float(d.get(y_col, d.get("count", 0))) for d in data]
                bar_colors = [CHART_PALETTE[i % len(CHART_PALETTE)] for i in range(len(labels))]
                ax.bar(range(len(labels)), values, color=bar_colors, edgecolor="none", width=0.65)
                ax.set_xticks(range(len(labels)))
                ax.set_xticklabels(labels, rotation=40, ha="right", fontsize=8, color=tick_color)

            elif chart_type == "line":
                x_vals = [str(d.get(x_col, "")) for d in data]
                y_vals = [float(d.get(y_col, 0)) for d in data]
                ax.plot(range(len(x_vals)), y_vals, color=CHART_PALETTE[0], linewidth=2.5,
                        marker="o", markersize=4, markerfacecolor=CHART_PALETTE[1])
                ax.fill_between(range(len(x_vals)), y_vals, alpha=0.1, color=CHART_PALETTE[0])
                step = max(1, len(x_vals) // 8)
                ax.set_xticks(range(0, len(x_vals), step))
                ax.set_xticklabels([x_vals[i] for i in range(0, len(x_vals), step)],
                                   rotation=40, ha="right", fontsize=8, color=tick_color)

            elif chart_type == "pie":
                labels = [str(d.get("name", d.get(x_col, "")))[:20] for d in data]
                values = [float(d.get("value", d.get("count", 0))) for d in data]
                pie_colors = CHART_PALETTE[:len(labels)]
                wedges, texts, autotexts = ax.pie(
                    values, labels=labels, colors=pie_colors, autopct="%1.1f%%",
                    textprops={"color": text_color, "fontsize": 9},
                    pctdistance=0.8, startangle=90, wedgeprops={"edgecolor": "white", "linewidth": 1.5}
                )
                ax.set_facecolor(fig.get_facecolor())

            elif chart_type in ("histogram", "hist"):
                labels = [str(d.get(x_col, "")) for d in data]
                values = [float(d.get("count", d.get(y_col, 0))) for d in data]
                ax.bar(range(len(labels)), values, color=CHART_PALETTE[1], edgecolor=CHART_PALETTE[0],
                       width=0.9, alpha=0.85)
                step = max(1, len(labels) // 10)
                ax.set_xticks(range(0, len(labels), step))
                ax.set_xticklabels([labels[i] for i in range(0, len(labels), step)],
                                   rotation=40, ha="right", fontsize=7, color=tick_color)

            elif chart_type == "scatter":
                x_vals = [float(d.get(x_col, 0)) for d in data]
                y_vals = [float(d.get(y_col, 0)) for d in data]
                ax.scatter(x_vals, y_vals, c=CHART_PALETTE[0], alpha=0.6,
                          edgecolors=CHART_PALETTE[1], s=30)

            elif chart_type == "area":
                x_vals = [str(d.get(x_col, "")) for d in data]
                y_vals = [float(d.get(y_col, 0)) for d in data]
                ax.fill_between(range(len(x_vals)), y_vals, alpha=0.3, color=CHART_PALETTE[0])
                ax.plot(range(len(x_vals)), y_vals, color=CHART_PALETTE[0], linewidth=2)
                step = max(1, len(x_vals) // 8)
                ax.set_xticks(range(0, len(x_vals), step))
                ax.set_xticklabels([x_vals[i] for i in range(0, len(x_vals), step)],
                                   rotation=40, ha="right", fontsize=8, color=tick_color)

        ax.set_title(title, color=text_color, fontsize=13, fontweight="bold", pad=14)
        ax.tick_params(colors=tick_color, labelsize=9)
        ax.grid(axis="y", alpha=0.3, color=grid_color, linestyle="--")
        for spine in ax.spines.values():
            spine.set_color(spine_color)
        ax.spines["top"].set_visible(False)
        ax.spines["right"].set_visible(False)

        x_label = chart_config.get("x_label", "")
        y_label = chart_config.get("y_label", "")
        if x_label:
            ax.set_xlabel(x_label, color=tick_color, fontsize=10, labelpad=8)
        if y_label:
            ax.set_ylabel(y_label, color=tick_color, fontsize=10, labelpad=8)

        plt.tight_layout()
        buf = io.BytesIO()
        fig.savefig(buf, format="png", dpi=180, bbox_inches="tight",
                    facecolor=fig.get_facecolor(), edgecolor="none")
        plt.close(fig)
        buf.seek(0)
        return buf.getvalue()

    # ══════════════════════════════════════════════════
    #  PDF REPORT
    # ══════════════════════════════════════════════════
    def _pdf_styles(self):
        """Create professional PDF paragraph styles."""
        styles = getSampleStyleSheet()
        custom = {}
        custom["cover_title"] = ParagraphStyle(
            "CoverTitle", parent=styles["Title"], fontSize=32,
            textColor=colors.HexColor(BRAND_PRIMARY), spaceAfter=8,
            alignment=TA_CENTER, fontName="Helvetica-Bold",
        )
        custom["cover_sub"] = ParagraphStyle(
            "CoverSub", parent=styles["Normal"], fontSize=14,
            textColor=colors.HexColor(BRAND_MUTED), spaceAfter=4,
            alignment=TA_CENTER,
        )
        custom["section"] = ParagraphStyle(
            "Section", parent=styles["Heading1"], fontSize=18,
            textColor=colors.HexColor(BRAND_PRIMARY), spaceAfter=10,
            spaceBefore=20, fontName="Helvetica-Bold",
            borderColor=colors.HexColor(BRAND_PRIMARY), borderWidth=0,
            borderPadding=0,
        )
        custom["subsection"] = ParagraphStyle(
            "SubSection", parent=styles["Heading2"], fontSize=13,
            textColor=colors.HexColor(BRAND_DARK), spaceAfter=6,
            spaceBefore=12, fontName="Helvetica-Bold",
        )
        custom["body"] = ParagraphStyle(
            "Body", parent=styles["Normal"], fontSize=10,
            textColor=colors.HexColor(BRAND_TEXT), spaceAfter=6,
            leading=15, alignment=TA_JUSTIFY,
        )
        custom["body_bold"] = ParagraphStyle(
            "BodyBold", parent=custom["body"], fontName="Helvetica-Bold",
        )
        custom["small"] = ParagraphStyle(
            "Small", parent=styles["Normal"], fontSize=8,
            textColor=colors.HexColor(BRAND_MUTED), spaceAfter=4,
        )
        custom["footer"] = ParagraphStyle(
            "Footer", parent=styles["Normal"], fontSize=8,
            textColor=colors.HexColor(BRAND_MUTED), alignment=TA_CENTER,
        )
        return custom

    def _add_pdf_header_footer(self, canvas, doc, title="Analytix AI Report"):
        """Add header line and footer to every page."""
        canvas.saveState()
        w, h = A4
        # Header line
        canvas.setStrokeColor(colors.HexColor(BRAND_PRIMARY))
        canvas.setLineWidth(2)
        canvas.line(1.5 * cm, h - 1 * cm, w - 1.5 * cm, h - 1 * cm)
        canvas.setFont("Helvetica", 7)
        canvas.setFillColor(colors.HexColor(BRAND_MUTED))
        canvas.drawString(1.5 * cm, h - 0.85 * cm, title)
        # Footer
        canvas.setFont("Helvetica", 7)
        canvas.drawCentredString(w / 2, 0.6 * cm, f"Page {doc.page}")
        canvas.drawRightString(w - 1.5 * cm, 0.6 * cm,
                               f"Generated: {datetime.now().strftime('%B %d, %Y')}")
        canvas.restoreState()

    def _make_section_divider(self):
        """Horizontal rule for section breaks."""
        return HRFlowable(
            width="100%", thickness=1,
            color=colors.HexColor(BRAND_BORDER),
            spaceBefore=6, spaceAfter=12,
        )

    def _build_data_overview(self, df: pd.DataFrame, s: dict) -> list:
        """Build a data overview section."""
        elements = []
        elements.append(Paragraph("Dataset Overview", s["section"]))
        elements.append(self._make_section_divider())

        # Summary stats table
        num_cols = len(df.select_dtypes(include=[np.number]).columns)
        cat_cols = len(df.select_dtypes(include=["object", "category"]).columns)
        missing_pct = round(df.isnull().sum().sum() / (len(df) * len(df.columns)) * 100, 1)

        overview_data = [
            ["Property", "Value", "Property", "Value"],
            ["Total Records", f"{len(df):,}", "Total Columns", str(len(df.columns))],
            ["Numeric Columns", str(num_cols), "Categorical Columns", str(cat_cols)],
            ["Missing Data", f"{missing_pct}%", "Duplicate Rows", str(df.duplicated().sum())],
        ]
        ow = 3.8 * cm
        tbl = Table(overview_data, colWidths=[ow, ow, ow, ow])
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(BRAND_PRIMARY)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, 0), 10),
            ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
            ("TOPPADDING", (0, 0), (-1, 0), 8),
            ("BACKGROUND", (0, 1), (0, -1), colors.HexColor("#f1f5f9")),
            ("BACKGROUND", (2, 1), (2, -1), colors.HexColor("#f1f5f9")),
            ("FONTNAME", (0, 1), (0, -1), "Helvetica-Bold"),
            ("FONTNAME", (2, 1), (2, -1), "Helvetica-Bold"),
            ("FONTSIZE", (0, 1), (-1, -1), 9),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor(BRAND_BORDER)),
            ("TOPPADDING", (0, 1), (-1, -1), 6),
            ("BOTTOMPADDING", (0, 1), (-1, -1), 6),
            ("LEFTPADDING", (0, 0), (-1, -1), 8),
        ]))
        elements.append(tbl)
        elements.append(Spacer(1, 12))

        # Column details table
        elements.append(Paragraph("Column Details", s["subsection"]))
        col_data = [["Column Name", "Data Type", "Non-Null Count", "Unique Values", "Sample Values"]]
        for col in df.columns[:20]:
            dtype = str(df[col].dtype)
            non_null = f"{df[col].count():,}"
            nunique = str(df[col].nunique())
            samples = ", ".join(str(v) for v in df[col].dropna().unique()[:3])
            if len(samples) > 35:
                samples = samples[:35] + "..."
            col_data.append([str(col)[:25], dtype, non_null, nunique, samples])

        cw = [4 * cm, 2.5 * cm, 2.8 * cm, 2.5 * cm, 5.5 * cm]
        col_tbl = Table(col_data, colWidths=cw)
        col_tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(BRAND_DARK)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, 0), 9),
            ("FONTSIZE", (0, 1), (-1, -1), 8),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8fafc")]),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor(BRAND_BORDER)),
            ("TOPPADDING", (0, 0), (-1, -1), 5),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 5),
            ("LEFTPADDING", (0, 0), (-1, -1), 6),
        ]))
        elements.append(col_tbl)
        return elements

    def _build_kpi_section(self, kpis: list, s: dict) -> list:
        """Build KPI section with a professional table."""
        elements = []
        elements.append(Paragraph("Key Performance Indicators", s["section"]))
        elements.append(self._make_section_divider())

        # KPI table: 2 columns — KPI name | Value
        kpi_data = [["KPI Metric", "Value", "Description"]]
        for kpi in kpis:
            kd = kpi if isinstance(kpi, dict) else kpi.dict() if hasattr(kpi, "dict") else kpi.model_dump()
            kpi_data.append([
                str(kd.get("label", "")),
                str(kd.get("value", "")),
                str(kd.get("description", ""))[:60],
            ])

        kpi_tbl = Table(kpi_data, colWidths=[5.5 * cm, 4 * cm, 8 * cm])
        kpi_tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(BRAND_PRIMARY)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, 0), 10),
            ("BOTTOMPADDING", (0, 0), (-1, 0), 8),
            ("TOPPADDING", (0, 0), (-1, 0), 8),
            ("FONTNAME", (0, 1), (0, -1), "Helvetica-Bold"),
            ("FONTNAME", (1, 1), (1, -1), "Helvetica-Bold"),
            ("TEXTCOLOR", (1, 1), (1, -1), colors.HexColor(BRAND_PRIMARY)),
            ("FONTSIZE", (0, 1), (-1, -1), 10),
            ("FONTSIZE", (2, 1), (2, -1), 8),
            ("TEXTCOLOR", (2, 1), (2, -1), colors.HexColor(BRAND_MUTED)),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8fafc")]),
            ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor(BRAND_BORDER)),
            ("TOPPADDING", (0, 1), (-1, -1), 7),
            ("BOTTOMPADDING", (0, 1), (-1, -1), 7),
            ("LEFTPADDING", (0, 0), (-1, -1), 8),
            ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
        ]))
        elements.append(kpi_tbl)
        return elements

    def _build_charts_section(self, charts: list, s: dict) -> list:
        """Build charts section — each chart on its own with description."""
        elements = []
        elements.append(PageBreak())
        elements.append(Paragraph("Visual Analysis", s["section"]))
        elements.append(self._make_section_divider())

        for i, chart in enumerate(charts):
            cd = chart if isinstance(chart, dict) else chart.dict() if hasattr(chart, "dict") else chart.model_dump()
            try:
                img_bytes = self.generate_chart_image(cd, width=7.5, height=4.5, style="light")
                img_buf = io.BytesIO(img_bytes)
                img = Image(img_buf, width=16 * cm, height=9.5 * cm)

                desc = cd.get("description", "")
                chart_block = [img]
                if desc:
                    chart_block.append(Spacer(1, 4))
                    chart_block.append(Paragraph(f"<i>{desc}</i>", s["small"]))
                chart_block.append(Spacer(1, 16))

                # Keep chart + description together, page break if needed
                elements.append(KeepTogether(chart_block))
            except Exception as e:
                logger.warning(f"Chart render failed: {e}")
                elements.append(Paragraph(f"Chart: {cd.get('title', 'N/A')} — could not render", s["body"]))
                elements.append(Spacer(1, 10))

        return elements

    def _build_insights_section(self, insights: list, s: dict) -> list:
        """Build AI insights section."""
        elements = []
        elements.append(PageBreak())
        elements.append(Paragraph("AI-Generated Insights", s["section"]))
        elements.append(self._make_section_divider())

        for i, ins in enumerate(insights):
            ind = ins if isinstance(ins, dict) else ins.dict() if hasattr(ins, "dict") else ins.model_dump()
            title = ind.get("title", f"Insight {i+1}")
            desc = ind.get("description", "")
            # Numbered insight
            elements.append(Paragraph(f"<b>{i+1}. {title}</b>", s["body_bold"]))
            if desc:
                elements.append(Paragraph(desc, s["body"]))
            elements.append(Spacer(1, 8))

        return elements

    def _build_data_sample(self, df: pd.DataFrame, s: dict) -> list:
        """Build a sample data table (first 15 rows)."""
        elements = []
        elements.append(PageBreak())
        elements.append(Paragraph("Data Sample (First 15 Records)", s["section"]))
        elements.append(self._make_section_divider())

        sample = df.head(15)
        # Limit to 8 columns max for readability
        display_cols = list(sample.columns[:8])
        header = [str(c)[:18] for c in display_cols]
        table_data = [header]
        for _, row in sample[display_cols].iterrows():
            table_data.append([str(v)[:20] if pd.notna(v) else "—" for v in row])

        n_cols = len(display_cols)
        col_w = min(2.5 * cm, 17 * cm / n_cols)
        tbl = Table(table_data, colWidths=[col_w] * n_cols)
        tbl.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor(BRAND_DARK)),
            ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
            ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
            ("FONTSIZE", (0, 0), (-1, 0), 7),
            ("FONTSIZE", (0, 1), (-1, -1), 7),
            ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f8fafc")]),
            ("GRID", (0, 0), (-1, -1), 0.4, colors.HexColor(BRAND_BORDER)),
            ("TOPPADDING", (0, 0), (-1, -1), 4),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 4),
            ("LEFTPADDING", (0, 0), (-1, -1), 4),
        ]))
        elements.append(tbl)
        if len(df.columns) > 8:
            elements.append(Spacer(1, 6))
            elements.append(Paragraph(
                f"<i>Showing {len(display_cols)} of {len(df.columns)} columns. "
                f"Full dataset contains {len(df):,} records.</i>", s["small"]))
        return elements

    def generate_pdf_report(self, dashboard_config: dict, insights: list[dict],
                             df: pd.DataFrame) -> bytes:
        """Generate a professional multi-page PDF report."""
        buf = io.BytesIO()
        doc = SimpleDocTemplate(
            buf, pagesize=A4,
            topMargin=1.3 * cm, bottomMargin=1.2 * cm,
            leftMargin=1.5 * cm, rightMargin=1.5 * cm,
        )
        s = self._pdf_styles()
        report_title = dashboard_config.get("title", "Data Analysis Report")
        elements = []

        # ── Cover Page ──
        elements.append(Spacer(1, 5 * cm))
        elements.append(Paragraph("ANALYTIX AI", s["cover_title"]))
        elements.append(Spacer(1, 8))
        elements.append(HRFlowable(
            width="60%", thickness=3, color=colors.HexColor(BRAND_PRIMARY),
            spaceBefore=0, spaceAfter=10,
        ))
        elements.append(Paragraph(report_title, ParagraphStyle(
            "ReportTitle", parent=s["cover_sub"], fontSize=18,
            textColor=colors.HexColor(BRAND_DARK), fontName="Helvetica-Bold",
        )))
        elements.append(Spacer(1, 12))
        elements.append(Paragraph(
            f"Records: {len(df):,}  |  Columns: {len(df.columns)}  |  "
            f"Generated: {datetime.now().strftime('%B %d, %Y at %I:%M %p')}",
            s["cover_sub"],
        ))
        elements.append(Spacer(1, 3 * cm))
        elements.append(Paragraph(
            "This report was automatically generated by Analytix AI using advanced "
            "machine learning analysis of your dataset. It contains key performance "
            "indicators, visual analysis, AI-generated insights, and a data overview.",
            ParagraphStyle("CoverDesc", parent=s["body"], alignment=TA_CENTER,
                          textColor=colors.HexColor(BRAND_MUTED)),
        ))
        elements.append(PageBreak())

        # ── Table of Contents ──
        elements.append(Paragraph("Table of Contents", s["section"]))
        elements.append(self._make_section_divider())
        toc_items = ["1. Dataset Overview", "2. Key Performance Indicators",
                      "3. Visual Analysis", "4. AI-Generated Insights",
                      "5. Data Sample"]
        for item in toc_items:
            elements.append(Paragraph(f"    {item}", ParagraphStyle(
                "TOC", parent=s["body"], fontSize=12, spaceAfter=10,
                textColor=colors.HexColor(BRAND_DARK),
            )))
        elements.append(PageBreak())

        # ── Sections ──
        elements.extend(self._build_data_overview(df, s))
        elements.extend(self._build_kpi_section(dashboard_config.get("kpis", []), s))
        elements.extend(self._build_charts_section(dashboard_config.get("charts", []), s))
        if insights:
            elements.extend(self._build_insights_section(insights, s))
        elements.extend(self._build_data_sample(df, s))

        # Build with header/footer
        doc.build(
            elements,
            onFirstPage=lambda c, d: self._add_pdf_header_footer(c, d, report_title),
            onLaterPages=lambda c, d: self._add_pdf_header_footer(c, d, report_title),
        )
        buf.seek(0)
        return buf.getvalue()

    # ══════════════════════════════════════════════════
    #  PPTX PRESENTATION
    # ══════════════════════════════════════════════════
    def _add_slide_bg(self, slide, dark=True):
        """Set slide background color."""
        slide.background.fill.solid()
        if dark:
            slide.background.fill.fore_color.rgb = RGBColor(0x0f, 0x17, 0x2a)
        else:
            slide.background.fill.fore_color.rgb = RGBColor(0xf8, 0xfa, 0xfc)

    def _add_slide_number(self, slide, num, total=None):
        """Add slide number in bottom-right."""
        text = f"{num}" if not total else f"{num} / {total}"
        txBox = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.4))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.color.rgb = RGBColor(0x64, 0x74, 0x8b)
        p.alignment = PP_ALIGN.RIGHT

    def _add_title_bar(self, slide, title_text, subtitle_text=""):
        """Add a branded title bar at the top of a slide."""
        # Title background bar
        shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), Inches(13.33), Inches(1.2)  # 1 = rectangle
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1e, 0x1b, 0x4b)
        shape.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.15), Inches(12), Inches(0.7))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.font.bold = True

        if subtitle_text:
            txBox2 = slide.shapes.add_textbox(Inches(0.6), Inches(0.75), Inches(12), Inches(0.4))
            tf2 = txBox2.text_frame
            p2 = tf2.paragraphs[0]
            p2.text = subtitle_text
            p2.font.size = Pt(12)
            p2.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)

    def _make_kpi_card(self, slide, x, y, w, h, label, value):
        """Draw a styled KPI card on a PPTX slide."""
        # Card background
        shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(0x1e, 0x29, 0x3b)
        shape.line.color.rgb = RGBColor(0x33, 0x41, 0x55)
        shape.line.width = Pt(1)

        # Label
        lbl_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.2), Inches(w - 0.3), Inches(0.4))
        tf = lbl_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = str(label)
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
        p.alignment = PP_ALIGN.CENTER

        # Value
        val_box = slide.shapes.add_textbox(
            Inches(x + 0.15), Inches(y + 0.7), Inches(w - 0.3), Inches(0.6))
        tf2 = val_box.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = str(value)
        p2.font.size = Pt(22)
        p2.font.color.rgb = RGBColor(0x63, 0x66, 0xf1)
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER

    def generate_pptx(self, dashboard_config: dict, insights: list[dict],
                       df: pd.DataFrame) -> bytes:
        """Generate a professional PPTX presentation."""
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        report_title = dashboard_config.get("title", "Data Analysis Report")
        kpis = dashboard_config.get("kpis", [])
        charts = dashboard_config.get("charts", [])
        slide_num = 0

        # Count total slides for numbering
        total_slides = 3 + (1 if kpis else 0) + len(charts) + (1 if insights else 0)

        # ── Slide 1: Title ──
        slide_num += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_bg(slide, dark=True)

        # Large brand name
        txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(11.33), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "ANALYTIX AI"
        p.font.size = Pt(54)
        p.font.color.rgb = RGBColor(0x63, 0x66, 0xf1)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Accent line
        line = slide.shapes.add_shape(1, Inches(4), Inches(3.1), Inches(5.33), Inches(0.06))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x63, 0x66, 0xf1)
        line.line.fill.background()

        # Report title
        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(3.4), Inches(11.33), Inches(1))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = report_title
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p2.alignment = PP_ALIGN.CENTER

        # Metadata
        txBox3 = slide.shapes.add_textbox(Inches(1), Inches(4.6), Inches(11.33), Inches(0.6))
        tf3 = txBox3.text_frame
        p3 = tf3.paragraphs[0]
        p3.text = (f"{len(df):,} Records  •  {len(df.columns)} Columns  •  "
                   f"{datetime.now().strftime('%B %d, %Y')}")
        p3.font.size = Pt(14)
        p3.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
        p3.alignment = PP_ALIGN.CENTER

        self._add_slide_number(slide, slide_num, total_slides)

        # ── Slide 2: Overview ──
        slide_num += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_bg(slide, dark=True)
        self._add_title_bar(slide, "Dataset Overview", f"{len(df):,} records across {len(df.columns)} columns")

        # Stats cards
        num_cols = len(df.select_dtypes(include=[np.number]).columns)
        cat_cols = len(df.select_dtypes(include=["object", "category"]).columns)
        missing_pct = round(df.isnull().sum().sum() / max(1, len(df) * len(df.columns)) * 100, 1)
        dupes = df.duplicated().sum()

        stats = [
            ("Total Records", f"{len(df):,}"),
            ("Columns", str(len(df.columns))),
            ("Numeric", str(num_cols)),
            ("Categorical", str(cat_cols)),
            ("Missing Data", f"{missing_pct}%"),
            ("Duplicates", str(dupes)),
        ]
        card_w, card_h, gap = 1.85, 1.4, 0.25
        start_x = (13.33 - (6 * card_w + 5 * gap)) / 2
        for i, (lbl, val) in enumerate(stats):
            x = start_x + i * (card_w + gap)
            self._make_kpi_card(slide, x, 1.8, card_w, card_h, lbl, val)

        # Column list
        col_box = slide.shapes.add_textbox(Inches(0.6), Inches(3.8), Inches(12), Inches(3.2))
        tf = col_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "Column Summary"
        p.font.size = Pt(16)
        p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
        p.font.bold = True

        cols_text = "  |  ".join([f"{c} ({df[c].dtype})" for c in df.columns[:15]])
        if len(df.columns) > 15:
            cols_text += f"  |  ... and {len(df.columns) - 15} more"
        p2 = tf.add_paragraph()
        p2.text = cols_text
        p2.font.size = Pt(10)
        p2.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
        p2.space_before = Pt(8)

        self._add_slide_number(slide, slide_num, total_slides)

        # ── Slide 3: KPIs ──
        if kpis:
            slide_num += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_bg(slide, dark=True)
            self._add_title_bar(slide, "Key Performance Indicators")

            cols_per_row = 4
            kpi_w, kpi_h = 2.8, 1.6
            kpi_gap = 0.25
            total_w = cols_per_row * kpi_w + (cols_per_row - 1) * kpi_gap
            kpi_start_x = (13.33 - total_w) / 2

            for i, kpi in enumerate(kpis[:8]):
                kd = kpi if isinstance(kpi, dict) else kpi.dict() if hasattr(kpi, "dict") else kpi.model_dump()
                row, col = i // 4, i % 4
                x = kpi_start_x + col * (kpi_w + kpi_gap)
                y = 1.8 + row * (kpi_h + 0.4)
                self._make_kpi_card(slide, x, y, kpi_w, kpi_h,
                                    kd.get("label", ""), kd.get("value", ""))

            self._add_slide_number(slide, slide_num, total_slides)

        # ── Chart Slides ──
        for chart in charts:
            slide_num += 1
            cd = chart if isinstance(chart, dict) else chart.dict() if hasattr(chart, "dict") else chart.model_dump()
            try:
                img_bytes = self.generate_chart_image(cd, width=11, height=5.5, style="dark")
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                self._add_slide_bg(slide, dark=True)

                chart_title = cd.get("title", "Chart")
                self._add_title_bar(slide, chart_title)

                img_buf = io.BytesIO(img_bytes)
                slide.shapes.add_picture(img_buf, Inches(1.2), Inches(1.4), Inches(10.9), Inches(5.4))

                # Description at bottom
                desc = cd.get("description", "")
                if desc:
                    desc_box = slide.shapes.add_textbox(Inches(1.2), Inches(6.9), Inches(10.9), Inches(0.5))
                    tf = desc_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = desc
                    p.font.size = Pt(10)
                    p.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
                    p.font.italic = True

                self._add_slide_number(slide, slide_num, total_slides)
            except Exception as e:
                logger.warning(f"PPTX chart render failed: {e}")

        # ── Insights Slide ──
        if insights:
            slide_num += 1
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            self._add_slide_bg(slide, dark=True)
            self._add_title_bar(slide, "Key Insights", "AI-generated analysis findings")

            ins_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.5))
            tf = ins_box.text_frame
            tf.word_wrap = True

            for i, ins in enumerate(insights[:8]):
                ind = ins if isinstance(ins, dict) else ins.dict() if hasattr(ins, "dict") else ins.model_dump()
                title = ind.get("title", f"Insight {i+1}")
                desc = ind.get("description", "")

                p = tf.add_paragraph()
                p.text = f"{i+1}. {title}"
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(0xff, 0xff, 0xff)
                p.font.bold = True
                p.space_before = Pt(10)

                if desc:
                    p2 = tf.add_paragraph()
                    p2.text = f"    {desc[:150]}"
                    p2.font.size = Pt(11)
                    p2.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
                    p2.space_before = Pt(2)

            self._add_slide_number(slide, slide_num, total_slides)

        # ── Thank You Slide ──
        slide_num += 1
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._add_slide_bg(slide, dark=True)

        txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.33), Inches(1.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "Thank You"
        p.font.size = Pt(48)
        p.font.color.rgb = RGBColor(0x63, 0x66, 0xf1)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        p2 = tf.add_paragraph()
        p2.text = "Report generated by Analytix AI"
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0x94, 0xa3, 0xb8)
        p2.alignment = PP_ALIGN.CENTER

        p3 = tf.add_paragraph()
        p3.text = datetime.now().strftime("%B %d, %Y")
        p3.font.size = Pt(12)
        p3.font.color.rgb = RGBColor(0x64, 0x74, 0x8b)
        p3.alignment = PP_ALIGN.CENTER

        self._add_slide_number(slide, slide_num, total_slides)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.getvalue()
